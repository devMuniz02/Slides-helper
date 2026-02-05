#!/usr/bin/env python3
"""
RAG (Retrieval-Augmented Generation) System for PowerPoint Presentations

This script provides a complete RAG system that:
1. Processes PowerPoint files to extract text and analyze images
2. Creates embeddings for each slide's content
3. Stores embeddings in a vector database
4. Provides an interactive chat interface for querying presentations

Usage:
    python rag_system.py --pptx path/to/presentation.pptx
    python rag_system.py --chat  # Start chat interface
    python rag_system.py --rebuild-db  # Rebuild vector database
"""

import argparse
import asyncio
import hashlib
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import uuid

import chromadb
from chromadb.config import Settings
import uvicorn
from contextlib import asynccontextmanager
from fastapi import FastAPI, WebSocket, WebSocketDisconnect, Request, UploadFile, File, Depends
from fastapi.responses import HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
import openai
from PIL import Image
from pptx import Presentation
import torch
from sentence_transformers import SentenceTransformer

# Import existing modules
from src.slide_processor import SlideProcessor, SlideContent
from src.vision_analyzer import VisionAnalyzer
from src.utils.config import config


class RAGSystem:
    """RAG system for PowerPoint presentations."""

    def __init__(self, db_path: str = "./rag_db"):
        """Initialize the RAG system."""
        self.db_path = Path(db_path)
        self.db_path.mkdir(exist_ok=True)

        # Initialize ChromaDB client
        self.chroma_client = chromadb.PersistentClient(
            path=str(self.db_path / "chroma_db"),
            settings=Settings(anonymized_telemetry=False)
        )

        # Initialize OpenAI client for generation and embeddings
        self.openai_client = openai.OpenAI(
            base_url=rag_config["lm_studio"]["base_url"],
            api_key="not-needed",  # LM Studio doesn't require API key
        )

        # Initialize embedding model (try LM Studio first, fallback to local)
        self.embedding_model = None
        self.use_lm_studio_embeddings = False
        self._init_embedding_model()

        # Initialize vision analyzer for image processing (optional)
        try:
            self.vision_analyzer = VisionAnalyzer()
            print("✓ Vision analyzer initialized")
        except Exception as e:
            print(f"⚠ Vision analyzer initialization failed: {e}")
            print("  Image analysis will be skipped")
            self.vision_analyzer = None

        # Initialize slide processor (will be created per PPTX file)
        self.slide_processor = None

        # Collections will be created per PPTX file
        self.collections = {}  # Dict to store collections by PPTX name

    def get_collection(self, pptx_name: str):
        """Get or create a collection for a specific PPTX file."""
        sanitized_name = sanitize_collection_name(f"slides_{pptx_name}")
        if sanitized_name not in self.collections:
            self.collections[sanitized_name] = self.chroma_client.get_or_create_collection(
                name=sanitized_name,
                metadata={"description": f"Embeddings for {pptx_name} PowerPoint content"}
            )
        return self.collections[sanitized_name]

    def _load_all_collections(self):
        """Load all existing collections from the database."""
        try:
            collections = self.chroma_client.list_collections()
            for collection in collections:
                if collection.name not in self.collections:
                    self.collections[collection.name] = collection
        except Exception as e:
            print(f"Error loading collections: {e}")

    def _init_embedding_model(self):
        """Initialize embedding model."""
        # Prefer local embeddings for better reliability and no LM Studio dependency
        print("Initializing local embeddings (sentence-transformers)...")
        
        try:
            self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
            self.use_lm_studio_embeddings = False
            print("✓ Using local sentence-transformers model (all-MiniLM-L6-v2)")
        except Exception as e:
            print(f"✗ Failed to initialize local embedding model: {e}")
            print("Falling back to LM Studio embeddings...")
            
            # Try LM Studio embeddings as fallback
            embedding_models_to_try = [
                rag_config["lm_studio"]["embedding_model_name"],
                "text-embedding-ada-002",
                "all-MiniLM-L6-v2",
                "text-embedding-3-small",
            ]
            
            for model_name in embedding_models_to_try:
                try:
                    print(f"  Testing LM Studio model: {model_name}")
                    response = self.openai_client.embeddings.create(
                        model=model_name,
                        input="test embedding"
                    )
                    if response.data and len(response.data) > 0:
                        self.use_lm_studio_embeddings = True
                        print(f"✓ Using LM Studio for embeddings with model: {model_name}")
                        return
                except Exception as e2:
                    print(f"  Model {model_name} failed: {str(e2)[:100]}...")
                    continue
            
            # If everything fails
            print("✗ No embedding models available!")
            print("Please ensure you have either:")
            print("1. Internet connection for sentence-transformers, OR")
            print("2. LM Studio running with an embedding model loaded")
            raise RuntimeError("No embedding models available")

    def get_embedding(self, text: str) -> List[float]:
        """Get embedding for text."""
        if self.use_lm_studio_embeddings:
            try:
                response = self.openai_client.embeddings.create(
                    model=rag_config["lm_studio"]["embedding_model_name"],
                    input=text
                )
                return response.data[0].embedding
            except Exception as e:
                print(f"LM Studio embedding failed, falling back to local: {e}")
                if self.embedding_model:
                    return self.embedding_model.encode(text).tolist()
                raise
        else:
            return self.embedding_model.encode(text).tolist()

    def process_pptx_file(self, pptx_path: str, force_reembed: bool = False) -> Dict[str, Any]:
        """
        Process a PowerPoint file and add to vector database.

        Args:
            pptx_path: Path to the PPTX file
            force_reprocess: If True, reprocess even if already exists

        Returns:
            Processing results
        """
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

        # Generate unique ID for this presentation
        pptx_hash = hashlib.md5(str(pptx_path).encode()).hexdigest()[:8]
        pptx_id = f"{pptx_path.stem}_{pptx_hash}"

        # Get collection for this PPTX
        collection = self.get_collection(pptx_path.stem)

        # Check if already processed
        existing_docs = collection.get(
            where={"pptx_id": pptx_id}
        )
        if existing_docs['ids'] and not force_reembed:
            print(f"✓ PPTX already processed: {pptx_path.name}")
            return {
                "status": "already_processed",
                "pptx_id": pptx_id,
                "slides_processed": len(existing_docs['ids'])
            }

        print(f"Processing PPTX: {pptx_path.name}")

        # Initialize slide processor for this PPTX
        slide_processor = SlideProcessor(str(pptx_path))

        # Extract slides
        slides = slide_processor.process_all_slides()

        # Process each slide
        processed_slides = []
        for slide in slides:
            slide_data = self._process_slide(slide, pptx_id, pptx_path.name, collection)
            processed_slides.append(slide_data)

        print(f"✓ Processed {len(processed_slides)} slides from {pptx_path.name}")

        return {
            "status": "processed",
            "pptx_id": pptx_id,
            "pptx_path": str(pptx_path),
            "slides_processed": len(processed_slides),
            "slides": processed_slides
        }

    def _process_slide(self, slide: SlideContent, pptx_id: str, pptx_filename: str, collection) -> Dict[str, Any]:
        """Process a single slide and add to vector database."""
        # Extract text content
        text_content = slide.get_full_text()

        # Analyze images if present
        image_descriptions = []
        if slide.images:
            print(f"  Analyzing {len(slide.images)} images on slide {slide.slide_number}...")
            for i, image in enumerate(slide.images):
                try:
                    if hasattr(self, 'vision_analyzer') and self.vision_analyzer:
                        description = self.vision_analyzer.analyze_slide_image(
                            image,
                            slide_text=text_content,
                            custom_prompt="Describe this image in detail, focusing on any charts, diagrams, or visual elements that convey information."
                        )
                        image_descriptions.append(f"Image {i+1}: {description}")
                    else:
                        print(f"    Vision analyzer not available, skipping image {i+1}")
                        image_descriptions.append(f"Image {i+1}: [Vision analysis not available]")
                except Exception as e:
                    print(f"Warning: Failed to analyze image {i+1} on slide {slide.slide_number}: {e}")
                    image_descriptions.append(f"Image {i+1}: [Image analysis failed]")

        # Combine all content
        combined_content = text_content
        if image_descriptions:
            combined_content += "\n\n" + "\n".join(image_descriptions)

        # Create embedding
        embedding = self.get_embedding(combined_content)

        # Prepare metadata
        metadata = {
            "pptx_id": pptx_id,
            "pptx_filename": pptx_filename,
            "slide_number": slide.slide_number,
            "has_text": bool(text_content.strip()),
            "has_images": bool(slide.images),
            "num_images": len(slide.images),
            "title": slide.title or "",
            "content_type": "text_and_images" if slide.images else "text_only"
        }

        # Add to vector database
        doc_id = f"{pptx_id}_slide_{slide.slide_number}"
        collection.add(
            embeddings=[embedding],
            documents=[combined_content],
            metadatas=[metadata],
            ids=[doc_id]
        )

        return {
            "slide_number": slide.slide_number,
            "doc_id": doc_id,
            "text_length": len(text_content),
            "images_analyzed": len(image_descriptions),
            "content_preview": combined_content[:200] + "..." if len(combined_content) > 200 else combined_content
        }

    def search_similar_content(self, query: str, top_k: int = 5, selected_pptx: List[str] = None) -> List[Dict[str, Any]]:
        """Search for similar content across selected PPTX collections."""
        # Ensure all collections are loaded
        self._load_all_collections()
        
        # Get query embedding
        query_embedding = self.get_embedding(query)

        # Determine which collections to search
        collections_to_search = {}
        if selected_pptx and len(selected_pptx) > 0:
            # Filter to selected collections (exact match)
            selected_lower = [s.lower() for s in selected_pptx]
            for collection_name, collection in self.collections.items():
                pptx_filename = collection_name.replace("slides_", "").lower()
                if pptx_filename in selected_lower:
                    collections_to_search[collection_name] = collection
        else:
            # Search all collections
            collections_to_search = self.collections

        # Search across selected collections
        all_results = []
        for collection_name, collection in collections_to_search.items():
            try:
                results = collection.query(
                    query_embeddings=[query_embedding],
                    n_results=top_k,
                    include=['documents', 'metadatas', 'distances']
                )

                # Add results with collection info
                for i, (doc, metadata, distance) in enumerate(zip(
                    results['documents'][0] if results['documents'] else [],
                    results['metadatas'][0] if results['metadatas'] else [],
                    results['distances'][0] if results['distances'] else []
                )):
                    all_results.append({
                        "content": doc,
                        "metadata": metadata,
                        "distance": distance,
                        "pptx_filename": metadata.get("pptx_filename", collection_name.replace("slides_", "")),
                        "slide_number": metadata.get("slide_number", 0),
                        "content_type": metadata.get("content_type", "unknown"),
                        "similarity_score": 1 - distance  # Convert distance to similarity
                    })
            except Exception as e:
                print(f"Error searching collection {collection_name}: {e}")
                continue
        
        # Sort by similarity (highest first)
        all_results.sort(key=lambda x: x["similarity_score"], reverse=True)
        return all_results[:top_k]

    def generate_answer(self, query: str, context_docs: List[Dict[str, Any]]) -> str:
        """Generate an answer using retrieved context."""
        # Prepare context
        context_parts = []
        for doc in context_docs:
            source_info = f"From {doc['pptx_filename']}, Slide {doc['slide_number']}"
            if doc['content_type'] == 'text_and_images':
                source_info += " (text and images)"
            elif doc['content_type'] == 'text_only':
                source_info += " (text)"

            context_parts.append(f"{source_info}:\n{doc['content']}")

        context = "\n\n".join(context_parts)

        # Create prompt
        prompt = f"""You are a helpful assistant that answers questions about PowerPoint presentations.

Use the following context from the presentations to answer the user's question. Each piece of context includes the source presentation filename and slide number.

Context:
{context}

Question: {query}

Instructions:
- Answer based on the provided context
- If the context doesn't contain enough information, say so
- Always cite your sources by mentioning the presentation filename and slide number
- Be concise but comprehensive
- If information comes from both text and images, mention that

Answer:"""

        # Generate response
        try:
            response = self.openai_client.chat.completions.create(
                model=config.lm_studio.model_name,
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that answers questions about PowerPoint presentations based on provided context."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=1000,
                temperature=0.3
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            return f"Error generating answer: {e}"

    def get_available_presentations(self) -> List[Dict[str, Any]]:
        """Get list of available presentations in the database."""
        try:
            # Get all documents
            results = self.collection.get(include=['metadatas'])

            # Group by presentation
            presentations = {}
            for metadata in results['metadatas']:
                pptx_id = metadata.get('pptx_id', 'unknown')
                pptx_filename = metadata.get('pptx_filename', 'Unknown')

                if pptx_id not in presentations:
                    presentations[pptx_id] = {
                        "pptx_id": pptx_id,
                        "filename": pptx_filename,
                        "slide_count": 0,
                        "has_images": False
                    }

                presentations[pptx_id]["slide_count"] += 1
                if metadata.get('has_images', False):
                    presentations[pptx_id]["has_images"] = True

            return list(presentations.values())
        except Exception as e:
            print(f"Error getting presentations: {e}")
            return []


# FastAPI app for RAG workflow interface
@asynccontextmanager
async def lifespan(app: FastAPI):
    """Handle application startup and shutdown."""
    # Startup
    global rag_system
    # Create uploads directory
    upload_dir = Path("./uploads")
    upload_dir.mkdir(exist_ok=True)
    # Don't initialize RAG system yet - wait for user configuration
    yield
    # Shutdown
    pass

app = FastAPI(title="PowerPoint RAG Workflow", lifespan=lifespan)
templates = Jinja2Templates(directory="templates")

import json
from pathlib import Path

def load_rag_config():
    """Load RAG configuration from rag_config.json file."""
    config_path = Path("config/rag_config.json")

    # Default configuration
    default_config = {
        "models": {
            "vision": "qwen/qwen2.5-vl-7b",
            "text_generation": "qwen2.5-7b-instruct",
            "embedding": "local"
        },
        "lm_studio": {
            "base_url": "http://localhost:1234/v1",
            "vision_model_name": "qwen/qwen2.5-vl-7b",
            "model_name": "qwen2.5-7b-instruct",
            "embedding_model_name": "text-embedding-ada-002"
        },
        "tts": {
            "engine": "edge",
            "voice": "es-ES-ElviraNeural"
        }
    }

    if config_path.exists():
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                user_config = json.load(f)
            # Merge user config with defaults
            config = default_config.copy()
            config.update(user_config)
            print(f"✓ Loaded configuration from {config_path}")
            return config
        except Exception as e:
            print(f"⚠ Error loading config file: {e}. Using defaults.")
            return default_config
    else:
        print(f"⚠ Config file {config_path} not found. Using defaults.")
        print("  You can create rag_config.json to customize model settings.")
        return default_config

# Load configuration
rag_config = load_rag_config()

# Global state
workflow_state = {
    "selected_pptx": [],  # List of selected PPTX file paths
    "file_status": {},    # Dict[file_path] -> {"extracted": bool, "embedded": bool, "json_path": str, "slides_count": int}
    "extracted_data": {}, # Dict[file_path] -> extracted data
    "chat_active": False,
    "models": rag_config["models"].copy()
}

# Global RAG system instance
rag_system = None

# Global logging flag
enable_logging = False

def check_existing_extraction(pptx_path: str) -> Optional[str]:
    """Check if a PPTX file has already been extracted and return the JSON path if found."""
    pptx_name = Path(pptx_path).stem
    json_path = Path("./rag_db") / f"{pptx_name}_extracted.json"

    if json_path.exists():
        return str(json_path)
    return None

def check_existing_embeddings(pptx_path: str) -> bool:
    """Check if a PPTX file has already been embedded in the vector database."""
    try:
        # Initialize RAG system if not already done
        global rag_system
        if not rag_system:
            rag_system = RAGSystem()
        
        pptx_hash = hashlib.md5(str(pptx_path).encode()).hexdigest()[:8]
        pptx_id = f"{Path(pptx_path).stem}_{pptx_hash}"

        # Get all collections from ChromaDB
        collections = rag_system.chroma_client.list_collections()
        
        # Check each collection for the pptx_id
        for collection in collections:
            try:
                existing_docs = collection.get(where={"pptx_id": pptx_id})
                if len(existing_docs['ids']) > 0:
                    return True
            except Exception as e:
                print(f"Error checking collection {collection.name}: {e}")
                continue

    except Exception as e:
        print(f"Error checking embeddings for {pptx_path}: {e}")
    
    return False

def sanitize_collection_name(name: str) -> str:
    """Sanitize a name to be valid for ChromaDB collection names."""
    import re
    # Replace spaces and invalid characters with underscores
    sanitized = re.sub(r'[^a-zA-Z0-9._-]', '_', name)
    # Ensure it starts and ends with alphanumeric
    sanitized = re.sub(r'^[^a-zA-Z0-9]+', '', sanitized)
    sanitized = re.sub(r'[^a-zA-Z0-9]+$', '', sanitized)
    # Ensure minimum length
    if len(sanitized) < 3:
        sanitized = f"pptx_{sanitized}" if sanitized else "pptx_collection"
    # Ensure maximum length
    if len(sanitized) > 512:
        sanitized = sanitized[:512]
    return sanitized

@app.get("/", response_class=HTMLResponse)
async def get_workflow_interface(request: Request):
    """Serve the main workflow interface."""
    global rag_system, enable_logging
    
    if enable_logging:
        print("🔄 Loading main workflow interface...")
    
    # Initialize RAG system if not already done
    if not rag_system:
        rag_system = RAGSystem()
    
    # Get available PPTX files
    pptx_files = []
    
    # Get existing collections to check for previously processed files
    try:
        existing_collections = rag_system.chroma_client.list_collections()
        collection_names = {col.name for col in existing_collections}
    except Exception as e:
        print(f"Error getting collections: {e}")
        collection_names = set()
    
    for file in Path(".").glob("*.pptx"):
        if enable_logging:
            print(f"📄 Loading file: {file.name}")
        
        file_path_str = str(file)
        status = workflow_state["file_status"].get(file_path_str, {})
        
        # Get accurate slide count and file size by reading the file
        try:
            presentation = Presentation(str(file))
            slides_count = len(presentation.slides)
        except Exception as e:
            print(f"Error reading {file_path_str}: {e}")
            slides_count = 0
        
        # Get file size in MB
        file_size_mb = file.stat().st_size / (1024 * 1024)
        
        # Check if file was previously processed (has a collection)
        sanitized_name = sanitize_collection_name(f"slides_{file.stem}")
        was_previously_processed = sanitized_name in collection_names
        
        # Consider file as "available" if it was previously processed OR currently fully processed
        is_available = was_previously_processed or (status.get("extracted", False) and status.get("embedded", False))
        
        # Get slide info for gallery display (only for extracted files)
        slides = []
        if status.get("extracted", False) or was_previously_processed:
            try:
                processor = SlideProcessor(file_path_str)
                presentation = Presentation(file_path_str)
                for i, slide in enumerate(presentation.slides):
                    title = processor.get_slide_title(slide)
                    slides.append({
                        "number": i + 1,
                        "title": title or f"Slide {i + 1}"
                    })
            except Exception as e:
                print(f"Error getting slide info for {file.name}: {e}")
                slides = []
        
        pptx_files.append({
            "name": file.name,
            "path": file_path_str,
            "size": file.stat().st_size,  # Keep original size for backward compatibility
            "size_mb": round(file_size_mb, 1),  # Add MB size
            "extracted": status.get("extracted", False),
            "embedded": status.get("embedded", False),
            "slides_count": slides_count,  # Use accurate slide count
            "slides": slides,
            "was_previously_processed": was_previously_processed,
            "is_available": is_available
        })
    
    web_config = config.get("web_interface", {})
    host = web_config.get("host", "localhost")
    port = web_config.get("port", 8000)
    
    if enable_logging:
        print(f"✅ Loaded {len(pptx_files)} PPTX files for workflow interface")
    
    return templates.TemplateResponse(
        "workflow.html",
        {
            "request": request,
            "pptx_files": pptx_files,
            "state": workflow_state,
            "host": host,
            "port": port
        }
    )

@app.get("/chat", response_class=HTMLResponse)
async def get_chat_interface(request: Request):
    """Serve the chat interface."""
    # Get available presentations from the database
    presentations = []
    try:
        # Query the database for available presentations
        results = rag_system.collection.get(include=["metadatas"])
        if results and results["metadatas"]:
            # Group by presentation filename
            pres_dict = {}
            for metadata in results["metadatas"]:
                filename = metadata.get("filename", "Unknown")
                if filename not in pres_dict:
                    pres_dict[filename] = {
                        "filename": filename,
                        "slide_count": 0,
                        "has_images": False
                    }
                pres_dict[filename]["slide_count"] = max(
                    pres_dict[filename]["slide_count"], 
                    metadata.get("slide_number", 0)
                )
                if metadata.get("content_type") == "image":
                    pres_dict[filename]["has_images"] = True
            
            presentations = list(pres_dict.values())
    except Exception as e:
        print(f"Warning: Could not load presentations: {e}")
    
    return templates.TemplateResponse(
        "chat.html",
        {
            "request": request,
            "presentations": presentations
        }
    )

@app.post("/upload-pptx")
async def upload_pptx(files: List[UploadFile] = File(...)):
    """Upload multiple PPTX files."""
    upload_dir = Path("./uploads")
    upload_dir.mkdir(exist_ok=True)

    uploaded_files = []
    for file in files:
        if not file.filename.lower().endswith(('.pptx', '.ppt')):
            continue

        file_path = upload_dir / file.filename
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)

        file_path_str = file_path.as_posix()

        # Check if already extracted
        existing_json = check_existing_extraction(file_path_str)
        extraction_status = "extracted" if existing_json else "not_extracted"

        # Update workflow state
        if file_path_str not in workflow_state["selected_pptx"]:
            workflow_state["selected_pptx"].append(file_path_str)

        if file_path_str not in workflow_state["file_status"]:
            workflow_state["file_status"][file_path_str] = {
                "extracted": extraction_status == "extracted",
                "embedded": False,
                "json_path": existing_json,
                "slides_count": 0
            }

        # Load existing data if available
        if existing_json:
            try:
                with open(existing_json, "r", encoding="utf-8") as f:
                    existing_data = json.load(f)
                    workflow_state["extracted_data"][file_path_str] = existing_data.get("slides", [])
                    workflow_state["file_status"][file_path_str]["slides_count"] = len(existing_data.get("slides", []))
            except Exception as e:
                print(f"Error loading existing extraction: {e}")

        uploaded_files.append({
            "filename": file.filename,
            "path": file_path_str,
            "extraction_status": extraction_status,
            "json_path": existing_json
        })

    workflow_state["current_step"] = "extract"
    return {"success": True, "files": uploaded_files}

@app.post("/select-pptx")
async def select_pptx(request: Request):
    """Select a PPTX file for processing."""
    data = await request.json()
    pptx_path = Path(data.get("pptx_path")).as_posix()
    
    if not pptx_path or not Path(pptx_path).exists():
        return {"success": False, "error": "Invalid PPTX file"}
    
    # Add to selected files if not already selected
    if pptx_path not in workflow_state["selected_pptx"]:
        workflow_state["selected_pptx"].append(pptx_path)

        # Initialize file status
        if pptx_path not in workflow_state["file_status"]:
            workflow_state["file_status"][pptx_path] = {
                "extracted": False,
                "embedded": False,
                "json_path": None,
                "slides_count": 0
            }

        # Check if already extracted
        json_path = check_existing_extraction(pptx_path)
        if json_path:
            workflow_state["file_status"][pptx_path]["extracted"] = True
            workflow_state["file_status"][pptx_path]["json_path"] = json_path
            # Load existing data
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    existing_data = json.load(f)
                    workflow_state["extracted_data"][pptx_path] = existing_data.get("slides", [])
                    workflow_state["file_status"][pptx_path]["slides_count"] = len(existing_data.get("slides", []))
            except Exception as e:
                print(f"Error loading existing extraction: {e}")

        # Check if already embedded
        if check_existing_embeddings(pptx_path):
            workflow_state["file_status"][pptx_path]["embedded"] = True

    return {
        "success": True,
        "pptx": Path(pptx_path).name,
        "total_selected": len(workflow_state["selected_pptx"]),
        "already_extracted": workflow_state["file_status"][pptx_path]["extracted"],
        "already_embedded": workflow_state["file_status"][pptx_path]["embedded"]
    }

@app.post("/deselect-pptx")
async def deselect_pptx(request: Request):
    """Deselect a PPTX file from processing."""
    data = await request.json()
    pptx_path = Path(data.get("pptx_path")).as_posix()
    
    print(f"Deselecting PPTX: {pptx_path}")
    print(f"Current selected files: {workflow_state['selected_pptx']}")
    
    if pptx_path in workflow_state["selected_pptx"]:
        workflow_state["selected_pptx"].remove(pptx_path)
        print(f"Successfully deselected: {pptx_path}")
        return {"success": True, "pptx": Path(pptx_path).name, "total_selected": len(workflow_state["selected_pptx"])}
    
    print(f"PPTX file not in selection: {pptx_path}")
    return {"success": False, "error": "PPTX file not in selection"}

@app.post("/update-models")
async def update_models(request: Request):
    """Update model configurations."""
    data = await request.json()
    workflow_state["models"].update(data)
    return {"success": True, "models": workflow_state["models"]}

@app.post("/delete-pptx-data")
async def delete_pptx_data(request: Request):
    """Delete extracted data and embeddings for a specific PPTX file."""
    data = await request.json()
    pptx_path = Path(data.get("pptx_path")).as_posix()
    
    if not pptx_path:
        return {"success": False, "error": "No PPTX path provided"}
    
    try:
        global rag_system
        if not rag_system:
            rag_system = RAGSystem()
        
        # Delete extracted JSON file
        json_path = check_existing_extraction(pptx_path)
        if json_path and Path(json_path).exists():
            Path(json_path).unlink()
            print(f"Deleted extracted JSON: {json_path}")
        
        # Delete collection
        pptx_stem = Path(pptx_path).stem
        collection_name = sanitize_collection_name(f"slides_{pptx_stem}")
        try:
            rag_system.chroma_client.delete_collection(collection_name)
            print(f"Deleted collection: {collection_name}")
            # Remove from loaded collections
            if collection_name in rag_system.collections:
                del rag_system.collections[collection_name]
        except Exception as e:
            print(f"Error deleting collection {collection_name}: {e}")
        
        # Update workflow state
        if pptx_path in workflow_state["file_status"]:
            workflow_state["file_status"][pptx_path]["extracted"] = False
            workflow_state["file_status"][pptx_path]["embedded"] = False
            workflow_state["file_status"][pptx_path]["json_path"] = None
            workflow_state["file_status"][pptx_path]["slides_count"] = 0
        
        if pptx_path in workflow_state["extracted_data"]:
            del workflow_state["extracted_data"][pptx_path]
        
        return {"success": True, "pptx": Path(pptx_path).name}
    
    except Exception as e:
        print(f"Error deleting data for {pptx_path}: {e}")
        return {"success": False, "error": str(e)}

@app.get("/get-workflow-state")
async def get_workflow_state():
    """Get current workflow state."""
    global rag_system, enable_logging
    
    if enable_logging:
        print("🔄 Loading workflow state...")
    
    # Initialize RAG system if not already done
    if not rag_system:
        rag_system = RAGSystem()
    
    # Get existing collections to check for previously processed files
    try:
        existing_collections = rag_system.chroma_client.list_collections()
        collection_names = {col.name for col in existing_collections}
    except Exception as e:
        print(f"Error getting collections: {e}")
        collection_names = set()
    
    # Build list of all PPTX files (available for selection)
    all_pptx_paths = set()
    
    # Check current directory
    for file in Path(".").glob("*.pptx"):
        all_pptx_paths.add(file)
    for file in Path(".").glob("*.ppt"):
        all_pptx_paths.add(file)
    
    # Check uploads directory
    uploads_dir = Path("./uploads")
    if uploads_dir.exists():
        for file in uploads_dir.glob("*.pptx"):
            all_pptx_paths.add(file)
        for file in uploads_dir.glob("*.ppt"):
            all_pptx_paths.add(file)
    
    if enable_logging:
        print(f"Found {len(all_pptx_paths)} unique PPTX files")
    
    all_pptx_files = []
    for file in sorted(all_pptx_paths):
        if enable_logging:
            print(f"📄 Processing file: {file}")
        if enable_logging:
            print(f"📄 Loading file: {file.name}")
        
        file_path_str = file.as_posix()
        status = workflow_state["file_status"].get(file_path_str, {})
        
        # Get accurate slide count and file size by reading the file
        try:
            presentation = Presentation(str(file))
            slides_count = len(presentation.slides)
        except Exception as e:
            print(f"Error reading {file_path_str}: {e}")
            slides_count = 0
        
        # Get file size in MB
        file_size_mb = file.stat().st_size / (1024 * 1024)
        
        # Check if file was previously processed (has a collection)
        sanitized_name = sanitize_collection_name(f"slides_{file.stem}")
        was_previously_processed = sanitized_name in collection_names
        
        # Consider file as "available" if it exists (all PPTX files are available for selection)
        all_pptx_files.append({
            "path": file_path_str,
            "filename": file.name,
            "slides_count": slides_count,
            "file_size_mb": round(file_size_mb, 1),
            "was_previously_processed": was_previously_processed,
            "currently_processed": status.get("extracted", False) and status.get("embedded", False),
            "is_rag_ready": was_previously_processed or (status.get("extracted", False) and status.get("embedded", False))
        })
    
    if enable_logging:
        print(f"✅ Loaded {len(all_pptx_files)} PPTX files")
    
    # Normalize paths to posix format for consistency
    normalized_selected = [Path(p).as_posix() for p in workflow_state["selected_pptx"]]
    normalized_file_status = {Path(k).as_posix(): v for k, v in workflow_state["file_status"].items()}
    normalized_extracted_data = {Path(k).as_posix(): v for k, v in workflow_state["extracted_data"].items()}
    
    workflow_state["selected_pptx"] = normalized_selected
    workflow_state["file_status"] = normalized_file_status
    workflow_state["extracted_data"] = normalized_extracted_data
    
    # Add all PPTX files to the response
    response = dict(workflow_state)
    response["all_pptx_files"] = all_pptx_files
    return response

@app.post("/extract-info")
async def extract_info(request: Request):
    """Extract text and image information from selected PPTX files."""
    data = await request.json()
    force_reextract = data.get("force_reextract", False)
    pptx_paths = data.get("pptx_paths", None)  # Optional: specific paths to process

    # Determine which files to process
    if pptx_paths:
        files_to_process = pptx_paths
    else:
        files_to_process = workflow_state["selected_pptx"]

    if not files_to_process or len(files_to_process) == 0:
        return {"success": False, "error": "No PPTX files selected"}

    try:
        processed_files = []
        total_slides = 0

        # Process each selected PPTX file
        for pptx_path in files_to_process:
            # Skip if already extracted and not forcing re-extraction
            if workflow_state["file_status"][pptx_path]["extracted"] and not force_reextract:
                print(f"Skipping already extracted PPTX: {pptx_path}")
                continue

            print(f"Processing PPTX: {pptx_path}")

            # Initialize slide processor for this PPTX
            slide_processor = SlideProcessor(pptx_path)
            slides = slide_processor.process_all_slides()

            extracted_data = []

            for slide in slides:
                slide_data = {
                    "pptx_filename": Path(pptx_path).name,
                    "slide_number": slide.slide_number,
                    "title": slide.title,
                    "text_content": slide.text_content,
                    "has_images": len(slide.images) > 0,
                    "image_count": len(slide.images),
                    "full_text": slide.get_full_text(),
                    "image_descriptions": []
                }

                # Analyze images if present and vision model is configured
                if slide.images and workflow_state["models"]["vision"]:
                    try:
                        # Initialize vision analyzer with selected model
                        vision_analyzer = VisionAnalyzer(
                            base_url=rag_config["lm_studio"]["base_url"],
                            api_key="not-needed",
                            model_name=workflow_state["models"]["vision"]
                        )

                        for i, image in enumerate(slide.images):
                            try:
                                description = vision_analyzer.analyze_slide_image(
                                    image,
                                    slide_text=slide_data["full_text"],
                                    custom_prompt="Describe this image in detail, focusing on any charts, diagrams, or visual elements that convey information."
                                )
                                slide_data["image_descriptions"].append({
                                    "image_index": i,
                                    "description": description
                                })
                            except Exception as e:
                                slide_data["image_descriptions"].append({
                                    "image_index": i,
                                    "description": f"[Analysis failed: {str(e)}]"
                                })
                    except Exception as e:
                        print(f"Vision analysis setup failed: {e}")

                extracted_data.append(slide_data)
                total_slides += 1

            # Save to separate JSON file for this PPTX
            pptx_name = Path(pptx_path).stem
            json_path = Path("./rag_db") / f"{pptx_name}_extracted.json"
            json_path.parent.mkdir(exist_ok=True)

            with open(json_path, "w", encoding="utf-8") as f:
                json.dump({
                    "pptx_path": pptx_path,
                    "extraction_date": str(datetime.now()),
                    "models_used": workflow_state["models"].copy(),
                    "total_slides": len(extracted_data),
                    "slides": extracted_data
                }, f, ensure_ascii=False, indent=2)

            # Update workflow state
            workflow_state["extracted_data"][pptx_path] = extracted_data
            workflow_state["file_status"][pptx_path]["extracted"] = True
            workflow_state["file_status"][pptx_path]["json_path"] = str(json_path)
            workflow_state["file_status"][pptx_path]["slides_count"] = len(extracted_data)

            processed_files.append({
                "pptx_path": pptx_path,
                "filename": Path(pptx_path).name,
                "slides_count": len(extracted_data),
                "json_path": str(json_path)
            })

        workflow_state["current_step"] = "embed"

        return {
            "success": True,
            "files_processed": len(processed_files),
            "total_slides": total_slides,
            "processed_files": processed_files
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@app.post("/create-embeddings")
async def create_embeddings(request: Request):
    """Create embeddings for extracted data."""
    data = await request.json()
    force_reembed = data.get("force_reembed", False)
    pptx_paths = data.get("pptx_paths", None)  # Optional: specific paths to process

    # Determine which files to process
    if pptx_paths:
        files_to_process = pptx_paths
    else:
        files_to_process = workflow_state["selected_pptx"]

    if not files_to_process or len(files_to_process) == 0:
        return {"success": False, "error": "No PPTX files selected"}

    try:
        global rag_system
        if not rag_system:
            rag_system = RAGSystem()

        processed_files = []
        total_slides_embedded = 0

        # Process each selected PPTX file
        for pptx_path in files_to_process:
            # Skip if already embedded and not forcing re-embedding
            if workflow_state["file_status"][pptx_path]["embedded"] and not force_reembed:
                print(f"Skipping already embedded PPTX: {pptx_path}")
                continue

            result = rag_system.process_pptx_file(pptx_path, force_reembed=force_reembed)
            if result["status"] == "processed":
                total_slides_embedded += result["slides_processed"]
                workflow_state["file_status"][pptx_path]["embedded"] = True

                processed_files.append({
                    "pptx_path": pptx_path,
                    "filename": Path(pptx_path).name,
                    "slides_embedded": result["slides_processed"]
                })
            elif result["status"] == "already_processed":
                processed_files.append({
                    "pptx_path": pptx_path,
                    "filename": Path(pptx_path).name,
                    "slides_embedded": result["slides_processed"],
                    "status": "already_processed"
                })

        workflow_state["current_step"] = "chat"

        return {
            "success": True,
            "files_processed": len(processed_files),
            "total_slides_embedded": total_slides_embedded,
            "processed_files": processed_files
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@app.post("/start-chat")
async def start_chat(request: Request):
    """Start the RAG chat session."""
    # Check if any files have been embedded
    any_embedded = any(
        file_status.get("embedded", False)
        for file_status in workflow_state["file_status"].values()
    )

    if not any_embedded:
        return {"success": False, "error": "No embeddings created yet. Please create embeddings for at least one PPTX file."}

    workflow_state["chat_active"] = True

    # Model will be used directly from workflow_state in stream_answer
    return {"success": True}

@app.websocket("/ws/chat")
async def chat_websocket(websocket: WebSocket):
    """Handle chat WebSocket connections with async token streaming."""
    await websocket.accept()

    try:
        # Ensure RAG system is initialized
        global rag_system
        if not rag_system:
            rag_system = RAGSystem()

        while True:
            # Receive message
            data = await websocket.receive_json()
            query = data.get("query", "").strip()

            if not query:
                continue

            # Search for relevant content (only in selected presentations)
            if workflow_state["selected_pptx"]:
                selected_filenames = [Path(pptx_path).stem for pptx_path in workflow_state["selected_pptx"]]
                search_results = rag_system.search_similar_content(query, top_k=3, selected_pptx=selected_filenames)
            else:
                search_results = rag_system.search_similar_content(query, top_k=3)

            if not search_results:
                await websocket.send_json({
                    "type": "answer",
                    "content": "No relevant information found in the presentations."
                })
                continue

            # Generate answer with streaming
            await stream_answer(websocket, query, search_results)

    except WebSocketDisconnect:
        print("WebSocket disconnected")
    except Exception as e:
        print(f"WebSocket error: {e}")
        await websocket.send_json({
            "type": "error",
            "content": f"An error occurred: {str(e)}"
        })


def highlight_answer_with_sources(answer: str, context_docs: List[Dict[str, Any]]) -> str:
    """Highlight parts of the answer that match content from different sources."""
    if not context_docs:
        return answer
    
    # Define colors for sources - one color per source in the list
    colors = ['#007bff', '#28a745', '#dc3545', '#ffc107', '#6f42c1', '#fd7e14', '#20c997', '#e83e8c', '#6c757d', '#17a2b8']
    
    highlighted_answer = answer
    
    # For each source, find matching phrases in the answer and highlight them
    for idx, doc in enumerate(context_docs):
        color = colors[idx % len(colors)]
        filename = doc.get('pptx_filename', '')
        slide_num = doc.get('slide_number', 0)
        content = doc.get('content', '').lower()
        
        # Skip if content is too short
        if len(content) < 20:
            continue
            
        answer_lower = highlighted_answer.lower()
        
        # Split source content into meaningful phrases (sentences)
        sentences = [s.strip() for s in content.split('.') if len(s.strip()) > 15]
        
        # Also try key phrases (3-4 word combinations)
        words = content.split()
        phrases = []
        for j in range(len(words) - 3):
            phrase = ' '.join(words[j:j+4])  # 4-word phrases
            if len(phrase) > 20:  # Substantial phrases only
                phrases.append(phrase)
        
        # Combine and limit to avoid over-highlighting
        text_chunks = (sentences + phrases[:2])[:3]  # Max 3 chunks per source
        
        # Highlight matching chunks in the answer
        for chunk in text_chunks:
            chunk_lower = chunk.lower()
            if chunk_lower in answer_lower and len(chunk) > 20:
                # Find all occurrences
                start_pos = 0
                while True:
                    start_idx = answer_lower.find(chunk_lower, start_pos)
                    if start_idx == -1:
                        break
                    
                    end_idx = start_idx + len(chunk_lower)
                    original_chunk = highlighted_answer[start_idx:end_idx]
                    
                    # Only highlight if not already highlighted
                    if '<span' not in original_chunk:
                        highlighted_chunk = f'<span style="color: {color}; font-weight: 500;">{original_chunk}</span>'
                        highlighted_answer = highlighted_answer[:start_idx] + highlighted_chunk + highlighted_answer[end_idx:]
                        
                        # Update positions for next search
                        answer_lower = highlighted_answer.lower()
                        start_pos = start_idx + len(highlighted_chunk)
                    else:
                        start_pos = end_idx
    
    # If no highlighting was applied, fall back to source citation highlighting
    if highlighted_answer == answer and 'Source:' in answer:
        # Split by source citations and highlight
        parts = re.split(r'(Source:[^.]+\.)', answer)
        result = []
        current_color = colors[0]
        
        for part in parts:
            if part.startswith('Source:'):
                # Extract source info and set color based on source index
                # For fallback, we need to map sources to colors
                result.append(f'<span style="color: {current_color}; font-weight: bold;">{part}</span>')
            else:
                if part.strip():
                    result.append(f'<span style="color: {current_color};">{part}</span>')
        
        return ''.join(result)
    
    return highlighted_answer

async def stream_answer(websocket: WebSocket, query: str, context_docs: List[Dict[str, Any]]):
    """Stream answer generation token by token."""
    # Prepare context
    context_parts = []
    for doc in context_docs:
        source_info = f"From {doc['pptx_filename']}, Slide {doc['slide_number']}"
        if doc['content_type'] == 'text_and_images':
            source_info += " (text and images)"
        elif doc['content_type'] == 'text_only':
            source_info += " (text)"

        context_parts.append(f"{source_info}:\n{doc['content']}")

    context = "\n\n".join(context_parts)

    # Create prompt
    prompt = f"""You are a helpful assistant that answers questions about PowerPoint presentations.

Use the following context from the presentations to answer the user's question. Each piece of context includes the source presentation filename and slide number.

Context:
{context}

Question: {query}

Instructions:
- Answer based on the provided context
- If the context doesn't contain enough information, say so
- Always cite your sources by mentioning the presentation filename and slide number
- Be concise but comprehensive
- If information comes from both text and images, mention that

Answer:"""

    try:
        # Start streaming response
        response = rag_system.openai_client.chat.completions.create(
            model=workflow_state["models"]["text_generation"],
            messages=[
                {"role": "system", "content": "You are a helpful assistant that answers questions about PowerPoint presentations based on provided context."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1000,
            temperature=0.3,
            stream=True  # Enable streaming
        )
        
        full_answer = ""
        for chunk in response:
            if chunk.choices[0].delta.content:
                token = chunk.choices[0].delta.content
                full_answer += token
                
                # Send token to client
                await websocket.send_json({
                    "type": "token",
                    "content": token
                })
        
        # Create highlighted answer
        highlighted_answer = highlight_answer_with_sources(full_answer, context_docs)
        
        # Send completion signal
        await websocket.send_json({
            "type": "answer_complete",
            "full_answer": full_answer,
            "highlighted_answer": highlighted_answer
        })

        # Send sources
        sources = []
        for result in context_docs:
            sources.append({
                "filename": result["pptx_filename"],
                "slide": result["slide_number"],
                "type": result["content_type"],
                "similarity": result["similarity_score"]
            })

        await websocket.send_json({
            "type": "sources",
            "sources": sources
        })

    except Exception as e:
        await websocket.send_json({
            "type": "error",
            "content": f"Error generating answer: {str(e)}"
        })


def create_html_template():
    """Create the HTML template for the chat interface."""
    template_dir = Path("templates")
    template_dir.mkdir(exist_ok=True)

    html_content = """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint RAG Chat</title>
    <style>
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 20px;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
        }
        .container {
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 10px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.3);
            overflow: hidden;
        }
        .header {
            background: #2c3e50;
            color: white;
            padding: 20px;
            text-align: center;
        }
        .presentations {
            padding: 20px;
            background: #f8f9fa;
            border-bottom: 1px solid #dee2e6;
        }
        .presentation-list {
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
        }
        .presentation-item {
            background: white;
            padding: 10px 15px;
            border-radius: 20px;
            border: 1px solid #dee2e6;
            font-size: 14px;
        }
        .chat-container {
            display: flex;
            height: 600px;
        }
        .chat-messages {
            flex: 1;
            padding: 20px;
            overflow-y: auto;
            background: #f8f9fa;
            border-right: 1px solid #dee2e6;
        }
        .message {
            margin-bottom: 15px;
            padding: 10px 15px;
            border-radius: 10px;
            max-width: 80%;
        }
        .message.user {
            background: #007bff;
            color: white;
            margin-left: auto;
            text-align: right;
        }
        .message.assistant {
            background: white;
            border: 1px solid #dee2e6;
        }
        .sources {
            background: #e9ecef;
            padding: 10px;
            border-radius: 5px;
            margin-top: 10px;
            font-size: 12px;
        }
        .source-item {
            margin-bottom: 5px;
        }
        .chat-input {
            padding: 20px;
            background: white;
            border-top: 1px solid #dee2e6;
        }
        .input-group {
            display: flex;
            gap: 10px;
        }
        #messageInput {
            flex: 1;
            padding: 12px;
            border: 1px solid #dee2e6;
            border-radius: 5px;
            font-size: 16px;
        }
        #sendButton {
            padding: 12px 20px;
            background: #28a745;
            color: white;
            border: none;
            border-radius: 5px;
            cursor: pointer;
            font-size: 16px;
        }
        #sendButton:hover {
            background: #218838;
        }
        #sendButton:disabled {
            background: #6c757d;
            cursor: not-allowed;
        }
        .typing {
            font-style: italic;
            color: #6c757d;
        }
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>PowerPoint RAG Chat</h1>
            <p>Ask questions about your PowerPoint presentations</p>
        </div>

        <div class="presentations">
            <h3>Available Presentations:</h3>
            <div class="presentation-list">
                {% for presentation in presentations %}
                <div class="presentation-item">
                    📄 {{ presentation.filename }} ({{ presentation.slide_count }} slides)
                    {% if presentation.has_images %}🖼️{% endif %}
                </div>
                {% endfor %}
            </div>
        </div>

        <div class="chat-container">
            <div class="chat-messages" id="messages">
                <div class="message assistant">
                    Hello! I'm here to help you find information in your PowerPoint presentations.
                    Ask me anything about the content, and I'll search through all available slides to provide relevant answers with source citations.
                </div>
            </div>
        </div>

        <div class="chat-input">
            <div class="input-group">
                <input type="text" id="messageInput" placeholder="Ask a question about your presentations..." />
                <button id="sendButton">Send</button>
            </div>
        </div>
    </div>

    <script>
        const ws = new WebSocket('ws://localhost:8000/ws/chat');
        const messagesDiv = document.getElementById('messages');
        const messageInput = document.getElementById('messageInput');
        const sendButton = document.getElementById('sendButton');

        let currentAssistantMessage = null;

        ws.onmessage = function(event) {
            const data = JSON.parse(event.data);

            if (data.type === 'answer') {
                if (currentAssistantMessage) {
                    currentAssistantMessage.innerHTML = data.content;
                } else {
                    addMessage('assistant', data.content);
                }
            } else if (data.type === 'sources') {
                if (currentAssistantMessage) {
                    let sourcesHtml = '<div class="sources"><strong>Sources:</strong><br>';
                    data.sources.forEach(source => {
                        sourcesHtml += `<div class="source-item">📄 ${source.filename}, Slide ${source.slide} (${source.type}) - Similarity: ${(source.similarity * 100).toFixed(1)}%</div>`;
                    });
                    sourcesHtml += '</div>';
                    currentAssistantMessage.innerHTML += sourcesHtml;
                }
            } else if (data.type === 'error') {
                addMessage('assistant', `❌ ${data.content}`);
            }
        };

        function addMessage(sender, content) {
            const messageDiv = document.createElement('div');
            messageDiv.className = `message ${sender}`;
            messageDiv.innerHTML = content;
            messagesDiv.appendChild(messageDiv);
            messagesDiv.scrollTop = messagesDiv.scrollHeight;

            if (sender === 'assistant') {
                currentAssistantMessage = messageDiv;
            }
        }

        function sendMessage() {
            const message = messageInput.value.trim();
            if (!message) return;

            // Add user message
            addMessage('user', message);

            // Clear input
            messageInput.value = '';

            // Disable send button
            sendButton.disabled = true;
            sendButton.textContent = 'Sending...';

            // Send to server
            ws.send(JSON.stringify({ query: message }));

            // Reset current assistant message
            currentAssistantMessage = null;
        }

        sendButton.addEventListener('click', sendMessage);
        messageInput.addEventListener('keypress', function(e) {
            if (e.key === 'Enter') {
                sendMessage();
            }
        });

        ws.onopen = function() {
            console.log('Connected to chat server');
        };

        ws.onclose = function() {
            console.log('Disconnected from chat server');
            addMessage('assistant', '❌ Connection lost. Please refresh the page.');
        };

        ws.onerror = function(error) {
            console.error('WebSocket error:', error);
            addMessage('assistant', '❌ Connection error. Please check if the server is running.');
        };

        // Re-enable send button when message is received
        const originalOnMessage = ws.onmessage;
        ws.onmessage = function(event) {
            originalOnMessage(event);
            sendButton.disabled = false;
            sendButton.textContent = 'Send';
        };
    </script>
</body>
</html>"""

    with open(template_dir / "chat.html", "w", encoding="utf-8") as f:
        f.write(html_content)


def main():
    """Main function."""
    parser = argparse.ArgumentParser(description="PowerPoint RAG System")
    parser.add_argument(
        "--pptx",
        type=str,
        help="Path to PowerPoint file to process"
    )
    parser.add_argument(
        "--chat",
        action="store_true",
        help="Start the chat interface"
    )
    parser.add_argument(
        "--rebuild-db",
        action="store_true",
        help="Rebuild the vector database from scratch"
    )
    parser.add_argument(
        "--db-path",
        type=str,
        default="./rag_db",
        help="Path to vector database directory"
    )
    parser.add_argument(
        "--force-reembed",
        action="store_true",
        help="Force re-embedding of slides even if already processed"
    )
    parser.add_argument(
        "--workflow",
        action="store_true",
        help="Start the workflow web interface (default)"
    )
    parser.add_argument(
        "--logging",
        action="store_true",
        help="Enable detailed logging of file loading operations"
    )

    args = parser.parse_args()

    # Set global logging flag
    global enable_logging
    enable_logging = args.logging

    # Initialize RAG system
    rag_system = RAGSystem(db_path=args.db_path)

    if args.rebuild_db:
        print("Rebuilding vector database...")
        # Clear existing database
        import shutil
        db_path = Path(args.db_path) / "chroma_db"
        if db_path.exists():
            shutil.rmtree(db_path)
        print("✓ Database cleared")

    if args.pptx:
        # Process PPTX file
        try:
            result = rag_system.process_pptx_file(args.pptx, force_reembed=args.force_reembed)
            print(f"✓ Processed: {result}")
        except Exception as e:
            print(f"✗ Error processing PPTX: {e}")
            sys.exit(1)

    # Get web interface config
    web_config = config.get("web_interface", {})
    host = web_config.get("host", "localhost")
    port = web_config.get("port", 8000)

    if args.chat:
        # Create HTML template
        create_html_template()

        # Start chat server (legacy)
        print(f"Starting legacy chat interface on http://{host}:{port}")
        print("Press Ctrl+C to stop")
        uvicorn.run(app, host=host, port=port)

    # If no arguments provided or workflow requested, start workflow interface
    if args.workflow or not any([args.pptx, args.chat, args.rebuild_db]):
        # Start workflow server (default)
        print(f"Starting workflow interface on http://{host}:{port}")
        print("Press Ctrl+C to stop")
        uvicorn.run(app, host=host, port=port)


if __name__ == "__main__":
    main()