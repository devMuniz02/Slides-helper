"""
Slide Processor Module

This module handles extraction of slides from PowerPoint files (.pptx),
including text content, images, and metadata.
"""
import io
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Dict, Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..utils.config import config
from ..utils.image_utils import resize_image


@dataclass
class SlideContent:
    """Represents the content of a single slide."""
    
    slide_number: int
    title: Optional[str]
    text_content: List[str]
    images: List[Image.Image]
    notes: Optional[str]
    metadata: Dict[str, Any]
    
    def get_full_text(self) -> str:
        """Get all text content as a single string."""
        parts = []
        if self.title:
            parts.append(f"Title: {self.title}")
        if self.text_content:
            parts.append("Content:\n" + "\n".join(self.text_content))
        if self.notes:
            parts.append(f"Notes: {self.notes}")
        return "\n\n".join(parts)


class SlideProcessor:
    """
    Processes PowerPoint presentations to extract content and images.
    
    This class provides modular functionality to:
    - Extract text from slides
    - Extract images from slides
    - Extract speaker notes
    - Convert slides to images
    """
    
    def __init__(self, pptx_path: str):
        """
        Initialize the slide processor.
        
        Args:
            pptx_path: Path to the PowerPoint file
        """
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        self.presentation = Presentation(str(self.pptx_path))
        self.slides_count = len(self.presentation.slides)
    
    def extract_text_from_shape(self, shape) -> List[str]:
        """
        Extract text from a shape.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            List of text strings found in the shape
        """
        texts = []
        
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        
        # Handle grouped shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                texts.extend(self.extract_text_from_shape(sub_shape))
        
        # Handle tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
        
        return texts
    
    def extract_images_from_slide(self, slide) -> List[Image.Image]:
        """
        Extract images from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            List of PIL Image objects
        """
        images = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_stream = io.BytesIO(shape.image.blob)
                    img = Image.open(image_stream)
                    
                    # Resize if needed based on quality setting
                    if config.processing.image_quality == "high":
                        img = resize_image(img, max_size=1920)
                    elif config.processing.image_quality == "medium":
                        img = resize_image(img, max_size=1280)
                    else:  # low
                        img = resize_image(img, max_size=854)
                    
                    images.append(img)
                except Exception as e:
                    print(f"Warning: Could not extract image from slide: {e}")
        
        return images
    
    def get_slide_title(self, slide) -> Optional[str]:
        """
        Extract the title from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Slide title or None if no title found
        """
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        return None
    
    def get_slide_notes(self, slide) -> Optional[str]:
        """
        Extract speaker notes from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Speaker notes or None if no notes found
        """
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if text_frame and text_frame.text.strip():
                return text_frame.text.strip()
        return None
    
    def process_slide(self, slide_index: int) -> SlideContent:
        """
        Process a single slide and extract all content.
        
        Args:
            slide_index: Zero-based index of the slide to process
            
        Returns:
            SlideContent object containing all extracted information
        """
        if slide_index >= self.slides_count:
            raise IndexError(f"Slide index {slide_index} out of range (total: {self.slides_count})")
        
        slide = self.presentation.slides[slide_index]
        
        # Extract title
        title = self.get_slide_title(slide)
        
        # Extract all text content
        text_content = []
        for shape in slide.shapes:
            if shape != slide.shapes.title:  # Skip title as we already extracted it
                texts = self.extract_text_from_shape(shape)
                text_content.extend(texts)
        
        # Extract images
        images = self.extract_images_from_slide(slide)
        
        # Extract notes
        notes = self.get_slide_notes(slide)
        
        # Build metadata
        metadata = {
            "slide_layout": slide.slide_layout.name if hasattr(slide, "slide_layout") else "Unknown",
            "has_images": len(images) > 0,
            "text_length": sum(len(t) for t in text_content),
        }
        
        return SlideContent(
            slide_number=slide_index + 1,  # 1-based numbering for humans
            title=title,
            text_content=text_content,
            images=images,
            notes=notes,
            metadata=metadata,
        )
    
    def process_all_slides(self) -> List[SlideContent]:
        """
        Process all slides in the presentation.
        
        Returns:
            List of SlideContent objects for all slides
        """
        return [self.process_slide(i) for i in range(self.slides_count)]
    
    def render_slide_to_image(self, slide_index: int, output_path: Optional[Path] = None) -> Image.Image:
        """
        Render a slide as an image.
        
        Note: This is a placeholder. Full slide rendering would require additional
        dependencies like python-pptx with COM automation or using external tools.
        For now, this extracts the first image or creates a simple representation.
        
        Args:
            slide_index: Zero-based index of the slide
            output_path: Optional path to save the image
            
        Returns:
            PIL Image object
        """
        slide_content = self.process_slide(slide_index)
        
        # If slide has images, return the first one
        if slide_content.images:
            img = slide_content.images[0]
        else:
            # Create a simple text-based representation
            # In a production system, you'd use a proper rendering library
            img = Image.new('RGB', (1920, 1080), color='white')
        
        if output_path:
            img.save(output_path)
        
        return img
    
    def get_presentation_info(self) -> Dict[str, Any]:
        """
        Get metadata about the presentation.
        
        Returns:
            Dictionary containing presentation information
        """
        return {
            "file_path": str(self.pptx_path),
            "file_name": self.pptx_path.name,
            "total_slides": self.slides_count,
            "has_notes": any(self.get_slide_notes(slide) for slide in self.presentation.slides),
        }
